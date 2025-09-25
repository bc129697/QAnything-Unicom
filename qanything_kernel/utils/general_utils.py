from sanic.request import Request
from sanic.exceptions import BadRequest
from qanything_kernel.utils.custom_log import debug_logger, embed_logger, rerank_logger
from qanything_kernel.configs.model_config import KB_SUFFIX, IMAGES_PROXY_URL
import inspect
import traceback
from urllib.parse import urlparse
import time
import os
import logging
import re
import requests
from functools import wraps
from datetime import datetime, timedelta
from bs4 import BeautifulSoup
from urllib.parse import urljoin
import html2text
import os
import docx2txt
import docx
import fitz  # PyMuPDF
import openpyxl
from pptx import Presentation
import email
import chardet
import mimetypes


__all__ = ['isURL', 'get_time', 'get_time_async', 'format_source_documents', 'format_source_documents_v1', 
           'format_source_documents_v2', 'safe_get', 'truncate_filename', "replace_image_references", 
           'shorten_data', 'read_files_with_extensions', 'validate_user_id', 'get_invalid_user_id_msg', 
           'clear_string', 'simplify_filename', 'string_bytes_length', 'correct_kb_id', 'clear_kb_id',
           'clear_string_is_equal', 'deduplicate_documents', 'fast_estimate_file_char_count',
           'check_user_id_and_user_info', 'get_table_infos', 'format_time_record', 'get_time_range',
           'html_to_markdown', "get_all_subpages", 'merge_source_documents', 'get_docx_toc']



def check_internet_connection(url='http://www.httpbin.org/get', timeout=5):
    try:
        response = requests.get(url, timeout=timeout)
        # 如果状态码为200，表示请求成功，说明有网络连接
        if response.status_code == 200:
            return True
    except requests.ConnectionError:
        # 捕获连接错误，意味着没有网络连接
        pass
    except requests.Timeout:
        # 请求超时
        print("Request timed out.")
    except Exception as e:
        # 其他异常情况
        print(f"An error occurred: {e}")
    return False



def get_invalid_user_id_msg(user_id):
    return "fail, Invalid user_id: {}. user_id 长度必须小于64，且必须只含有字母，数字和下划线且字母开头".format(user_id)


def isURL(string):
    result = urlparse(string)
    return result.scheme != '' and result.netloc != ''


def format_source_documents(ori_source_documents):
    source_documents = []
    for inum, doc in enumerate(ori_source_documents):
        source_info = {'file_id': doc.metadata.get('file_id', ''),
                       'file_name': doc.metadata.get('file_name', ''),
                       'content': doc.page_content,
                       'retrieval_query': doc.metadata.get('retrieval_query', ''),
                       # 'kernel': doc.metadata['kernel'],
                       'file_url': doc.metadata.get('file_url', ''),
                       'score': str(doc.metadata['score']),
                       'embed_version': doc.metadata.get('embed_version', ''),
                       'nos_keys': doc.metadata.get('nos_keys', ''),
                       'doc_id': doc.metadata.get('doc_id', ''),
                       'retrieval_source': doc.metadata.get('retrieval_source', ''),
                       'headers': doc.metadata.get('headers', {}),
                       'page_id': doc.metadata.get('page_id', 0),
                       }
        source_documents.append(source_info)
    return source_documents

def format_source_documents_v1(ori_source_documents):
    source_documents = []
    source_web_documents = []
    for inum, doc in enumerate(ori_source_documents):
        source_info = {
            "page_content": doc.page_content,
            "metadata": {
                       'kb_id': doc.metadata.get('kb_id', ''),
                       'file_id': doc.metadata.get('file_id', ''),
                       'file_name': doc.metadata.get('file_name', ''),
                       'retrieval_query': doc.metadata.get('retrieval_query', ''),
                       'file_url': doc.metadata.get('file_url', ''),
                       'score': float(doc.metadata['score']),
                       'doc_id': doc.metadata.get('doc_id', ''),
                       'retrieval_source': doc.metadata.get('retrieval_source', ''),
                       'headers': doc.metadata.get('headers', {}),
                       'page_id': doc.metadata.get('page_id', 0)
                    }
        }
        if doc.metadata.get('file_name', '').endswith('.faq'):
            source_info["metadata"]["faq_dict"] = doc.metadata.get("faq_dict", {})
        if doc.metadata.get('file_name', '').endswith('.web'):
            source_web_documents.append(source_info)
        else:
            source_documents.append(source_info)
    return source_documents, source_web_documents

def format_source_documents_v2(ori_source_documents):
    source_doc_documents = []
    source_qa_documents = []
    source_web_documents = []
    for inum, doc in enumerate(ori_source_documents):
        source_info = {
            "page_content": doc.page_content,
            "metadata": {
                       'kb_id': doc.metadata.get('kb_id', ''),
                       'file_id': doc.metadata.get('file_id', ''),
                       'file_name': doc.metadata.get('file_name', ''),
                       'retrieval_query': doc.metadata.get('retrieval_query', ''),
                       'file_url': doc.metadata.get('file_url', ''),
                       'score': float(doc.metadata['score']),
                       'doc_id': doc.metadata.get('doc_id', ''),
                       'retrieval_source': doc.metadata.get('retrieval_source', ''),
                       'headers': doc.metadata.get('headers', {}),
                       'page_id': doc.metadata.get('page_id', 0)
            }
        }
        if doc.metadata.get('file_name', '').endswith('.faq'):
            source_info["metadata"]["faq_dict"] = doc.metadata.get("faq_dict", {})
            source_qa_documents.append(source_info)
        elif doc.metadata.get('file_name', '').endswith('.web'):
            source_web_documents.append(source_info)
        else:
            source_doc_documents.append(source_info)
    return source_doc_documents, source_qa_documents, source_web_documents


def merge_source_documents(ori_source_documents):
    # 按file_id分组并记录知识库信息
    file_groups = {}
    for doc in ori_source_documents:
        meta = doc["metadata"]
        file_id = meta["file_id"]
        
        if file_id not in file_groups:
            file_groups[file_id] = {
                "metadata": meta,
                "slices": [],
                "score": []
            }
        
        # 解析切片ID（假设doc_id格式为"fileid_sliceid"）
        slice_id = int(meta["doc_id"].split("_")[-1])
        file_groups[file_id]["slices"].append((slice_id, doc["page_content"]))
        file_groups[file_id]["score"].append(meta["score"])
    
    merge_source_documents = []
    for index, (file_id, group) in enumerate(file_groups.items(), 1):
        # 排序切片内容
        sorted_slices = sorted(group["slices"], key=lambda x: x[0])
        
        # 构建文件块
        file_block = [
            # f"***参考信息[{index}]：*** （***参考信息来源：*** 文件名: {group['file_name']}; 知识库名: {group['kb_name']}）"
        ]
        
        if len(sorted_slices) > 1:
            # 添加切片内容
            file_block.append("该参考内容由多个片段组成，片段内容如下：")
            for slice_id, content in sorted_slices:
                # 清理内容中的多余换行
                cleaned_content = content.strip().replace('\n\n', '\n')
                file_block.append(f"***片段[{slice_id}]：*** {cleaned_content}")
        else:
            # 如果只有一个切片，直接添加内容
            cleaned_content = sorted_slices[0][1].strip().replace('\n\n', '\n')
            file_block.append(cleaned_content)
        
        page_content = "\n\n".join(file_block)
        source_info = {
            "page_content": page_content,
            "metadata": {
                'kb_id': group["metadata"].get('kb_id', ''),
                'file_id': group["metadata"].get('file_id', ''),
                'file_name': group["metadata"].get('file_name', ''),
                'retrieval_query': group["metadata"].get('retrieval_query', ''),
                'file_url': group["metadata"].get('file_url', ''),
                'score': max(group["score"]),
                'retrieval_source': group["metadata"].get('retrieval_source', ''),
                'headers': group["metadata"].get('headers', {}),
                'page_id': group["metadata"].get('page_id', 0),
                }
        }
        if group["metadata"].get('faq_dict', None):
            source_info["metadata"]["faq_dict"] = group["metadata"]["faq_dict"]
        merge_source_documents.append(source_info)

    return merge_source_documents
        


def format_time_record(time_record):
    token_usage = {}
    time_usage = {}
    for k, v in time_record.items():
        if 'tokens' in k:
            token_usage[k] = round(v)
        else:
            time_usage[k] = round(v, 2)
    if 'rewrite_prompt_tokens' in token_usage:
        if 'prompt_tokens' in token_usage:
            token_usage['prompt_tokens'] += token_usage['rewrite_prompt_tokens']
        if 'total_tokens' in token_usage:
            token_usage['total_tokens'] += token_usage['rewrite_prompt_tokens']
    if 'rewrite_completion_tokens' in token_usage:
        if 'completion_tokens' in token_usage:
            token_usage['completion_tokens'] += token_usage['rewrite_completion_tokens']
        if 'total_tokens' in token_usage:
            token_usage['total_tokens'] += token_usage['rewrite_completion_tokens']
    return {"time_usage": time_usage, "token_usage": token_usage}


def safe_get(req: Request, attr: str, default=None):
    try:
        if attr in req.form:
            return req.form.getlist(attr)[0]
        if attr in req.args:
            return req.args[attr]
        if attr in req.json:
            return req.json[attr]
        # if value := req.form.get(attr):
        #     return value
        # if value := req.args.get(attr):
        #     return value
        # """req.json执行时不校验content-type，body字段可能不能被正确解析为json"""
        # if value := req.json.get(attr):
        #     return value
    except BadRequest:
        logging.warning(f"missing {attr} in request")
    except Exception as e:
        logging.warning(f"get {attr} from request failed:")
        logging.warning(traceback.format_exc())
    return default


def truncate_filename(filename, max_length=200):
    # 获取文件名后缀
    file_ext = os.path.splitext(filename)[1]

    # 获取不带后缀的文件名
    file_name_no_ext = os.path.splitext(filename)[0]

    # 计算文件名长度，注意中文字符
    filename_length = len(filename.encode('utf-8'))

    # 如果文件名长度超过最大长度限制
    if filename_length > max_length:
        debug_logger.warning("文件名长度超过最大长度限制，将截取文件名")
        # 生成一个时间戳标记
        timestamp = str(int(time.time()))
        # 截取文件名
        while filename_length > max_length:
            file_name_no_ext = file_name_no_ext[:-4]
            new_filename = file_name_no_ext + "_" + timestamp + file_ext
            filename_length = len(new_filename.encode('utf-8'))
    else:
        new_filename = filename

    return new_filename


# 同步执行环境下的耗时统计装饰器
def get_time(func):
    def get_time_inner(*arg, **kwargs):
        s_time = time.time()
        res = func(*arg, **kwargs)
        e_time = time.time()
        if 'embed' in func.__name__:
            embed_logger.info('函数 {} 执行耗时: {:.2f} 秒'.format(func.__name__, e_time - s_time))
        elif 'rerank' in func.__name__:
            rerank_logger.info('函数 {} 执行耗时: {:.2f} 秒'.format(func.__name__, e_time - s_time))
        else:
            debug_logger.info('函数 {} 执行耗时: {:.2f} 毫秒'.format(func.__name__, (e_time - s_time) * 1000))
        return res

    return get_time_inner


# 异步执行环境下的耗时统计装饰器
def get_time_async(func):
    @wraps(func)
    async def get_time_async_inner(*args, **kwargs):
        s_time = time.perf_counter()
        res = await func(*args, **kwargs)  # 注意这里使用 await 来调用异步函数
        e_time = time.perf_counter()
        if 'embed' in func.__name__:
            embed_logger.info('函数 {} 执行耗时: {:.2f} 秒'.format(func.__name__, e_time - s_time))
        elif 'rerank' in func.__name__:
            rerank_logger.info('函数 {} 执行耗时: {:.2f} 秒'.format(func.__name__, e_time - s_time))
        else:
            debug_logger.info('函数 {} 执行耗时: {:.2f} 毫秒'.format(func.__name__, (e_time - s_time) * 1000))
        return res

    return get_time_async_inner


def read_files_with_extensions():
    # 获取当前脚本文件的路径
    current_file = os.path.abspath(__file__)

    # 获取当前脚本文件所在的目录
    current_dir = os.path.dirname(current_file)

    # 获取项目根目录
    project_dir = os.path.dirname(os.path.dirname(current_dir))

    directory = project_dir + '/data'

    extensions = ['.md', '.txt', '.pdf', '.jpg', '.docx', '.xlsx', '.eml', '.csv', 'pptx', 'jpeg', 'png']

    files = []
    for root, dirs, files_list in os.walk(directory):
        for file in files_list:
            if file.endswith(tuple(extensions)):
                file_path = os.path.join(root, file)
                with open(file_path, 'rb') as f:
                    file_content = f.read()
                    mime_type, _ = mimetypes.guess_type(file_path)
                    if mime_type is None:
                        mime_type = 'application/octet-stream'
                    # 模拟 req.files.getlist('files') 返回的对象
                    file_obj = type('FileStorage', (object,), {
                        'name': file,
                        'type': mime_type,
                        'body': file_content
                    })()
                    files.append(file_obj)
    return files


def validate_user_id(user_id):
    if len(user_id) > 64:
        return False
    # 定义正则表达式模式
    pattern = r'^[A-Za-z][A-Za-z0-9_]*$'
    # 检查是否匹配
    if isinstance(user_id, str) and re.match(pattern, user_id):
        return True
    else:
        return False


def shorten_data(data):
    # copy data，不要修改原始数据
    data = data.copy()
    try:
        for k, v in data.items():
            if len(str(v)) > 100:
                data[k] = str(v)[:100] + '...'
    except Exception as e:
        debug_logger.error('shorten_data error:', traceback.format_exc())
    return data


def cur_func_name():
    return inspect.currentframe().f_back.f_code.co_name



def sent_tokenize(x):
    #  sents_temp = re.split('(：|:|,|，|。|！|\!|\.|？|\?)', x)
    sents_temp = re.split('(。|！|\!|\.|？|\?)', x)
    sents = []
    for i in range(len(sents_temp) // 2):
        sent = sents_temp[2 * i] + sents_temp[2 * i + 1]
        sents.append(sent)
    return sents


def clear_string(str):
    # 只保留中文、英文、数字
    str = re.sub(r"[^\u4e00-\u9fa5a-zA-Z0-9]", "", str)
    return str


def simplify_filename(filename, max_length=40):
    if len(filename) <= max_length:
        # 如果文件名长度小于等于最大长度，直接返回原文件名
        return filename

    # 分离文件的基本名和扩展名
    name, extension = filename.rsplit('.', 1)
    extension = '.' + extension  # 将点添加回扩展名

    # 计算头部和尾部的保留长度
    part_length = (max_length - len(extension) - 1) // 2  # 减去扩展名长度和破折号的长度
    end_start = -part_length if part_length else None

    # 构建新的简化文件名
    simplified_name = f"{name[:part_length]}-{name[end_start:]}" if part_length else name[:max_length - 1]

    return f"{simplified_name}{extension}"


# 对比两个字符串，只保留字母数字和中文，返回是否一致
def clear_string_is_equal(str1, str2):
    str1 = clear_string(str1)
    str2 = clear_string(str2)
    return str1 == str2


def correct_kb_id(kb_id):
    if not kb_id:
        return kb_id
    # 如果kb_id末尾不是KB_SUFFIX,则加上
    if KB_SUFFIX not in kb_id:
        if kb_id.endswith('_FAQ'):  # KBc86eaa3f278f4ef9908780e8e558c6eb_FAQ
            return kb_id.split('_FAQ')[0] + KB_SUFFIX + '_FAQ'
        else:  # KBc86eaa3f278f4ef9908780e8e558c6eb
            return kb_id + KB_SUFFIX
    else:
        return kb_id


def clear_kb_id(kb_id):
    return kb_id.replace(KB_SUFFIX, '')


def string_bytes_length(string):
    return len(string.encode('utf-8'))


def check_user_id_and_user_info(user_id, user_info):
    if user_id is None or user_info is None:
        msg = "fail, user_id 或 user_info 为 None"
        return False, msg
    # if not validate_user_id(user_id):
    #     msg = get_invalid_user_id_msg(user_id)
    #     return False, msg
    # if not user_info.isdigit():
    #     msg = "fail, user_info 必须是纯数字"
    #     return False, msg
    return True, 'userid check success.'



def get_table_infos(markdown_str):
    lines = markdown_str.split('\n')
    if len(lines) < 2:
        return None
    head_line = None
    end_line = None
    for i in range(len(lines) - 1):
        if '|' in lines[i] and '|' in lines[i + 1]:
            separator_line = lines[i + 1].strip()
            if separator_line.startswith('|') and separator_line.endswith('|'):
                separator_parts = separator_line[1:-1].split('|')
                if all(part.strip().startswith('-') and len(part.strip()) >= 3 for part in separator_parts):
                    head_line = i
                    break
    for i in range(len(lines)):
        if '|' in lines[i]:
            separator_line = lines[i].strip()
            if separator_line.startswith('|') and separator_line.endswith('|'):
                end_line = i
    if head_line is None or end_line is None:
        return None
    return {"head_line": head_line, "end_line": end_line, "head": lines[head_line] + '\n' + lines[head_line + 1],
            "lines": lines}


def get_time_range(time_start=None, time_end=None, default_days=30):
    """
    获取时间范围。如果给定的时间范围不完整，将使用默认值（最近30天）。

    :param time_start: 起始时间，格式为 "YYYY-MM-DD" 或 "YYYY-MM-DD HH:MM:SS"
    :param time_end: 结束时间，格式为 "YYYY-MM-DD" 或 "YYYY-MM-DD HH:MM:SS"
    :param default_days: 如果未提供时间范围，默认的天数范围
    :return: 包含起始时间和结束时间的元组，格式为 ("YYYY-MM-DD HH:MM:SS", "YYYY-MM-DD HH:MM:SS")
    """

    date_pattern = r'^\d{4}-\d{2}-\d{2}$'
    now = datetime.now()

    # 验证 time_start 格式
    if time_start:
        if not re.match(date_pattern, time_start):
            return None
        if len(time_start) == 10:
            time_start = time_start + " 00:00:00"
    else:
        time_start = (now - timedelta(days=default_days)).strftime("%Y-%m-%d 00:00:00")

    # 验证 time_end 格式
    if time_end:
        if not re.match(date_pattern, time_end):
            return None
        if len(time_end) == 10:
            time_end = time_end + " 23:59:59"
    else:
        time_end = now.strftime("%Y-%m-%d 23:59:59")

    return (time_start, time_end)


def deduplicate_documents(source_docs):
    unique_docs = set()
    deduplicated_docs = []
    for doc in source_docs:
        if doc.page_content not in unique_docs:
            unique_docs.add(doc.page_content)
            deduplicated_docs.append(doc)
    return deduplicated_docs


def get_all_subpages(url):
    headers = {
        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
    }

    try:
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()
    except requests.RequestException as e:
        print(f"Error fetching {url}: {e}")
        return []

    soup = BeautifulSoup(response.text, 'html.parser')
    links = soup.find_all('a', href=True)
    subpages = set()

    for link in links:
        href = link['href']
        full_url = urljoin(url, href)
        subpages.add(full_url)

    return list(subpages)


def html_to_markdown(html_content):
    # 创建HTML到文本转换器
    h = html2text.HTML2Text()

    # 配置转换器
    h.ignore_images = True
    h.ignore_emphasis = True
    h.ignore_links = True
    h.body_width = 0  # 禁用换行
    h.tables = True  # 保留表格

    # 转换HTML到Markdown
    markdown = h.handle(html_content)

    # 删除所有图片标记
    markdown = re.sub(r'!\[.*?\]\(.*?\)', '', markdown)

    # 删除所有链接标记，保留文字
    markdown = re.sub(r'\[([^\]]*)\]\(.*?\)', r'\1', markdown)

    # 删除多余的空行，但保留表格结构
    # markdown = re.sub(r'(\n\s*){3,}', '\n\n', markdown)

    # 删除行首的特殊字符（如*、-等），但保留表格的|符号
    # markdown = re.sub(r'^(?!\|)\s*[-*]\s+', '', markdown, flags=re.MULTILINE)

    return markdown.strip()


def fast_estimate_file_char_count(file_path):
    """
    快速估算文件的字符数，如果超过max_chars则返回False，否则返回True
    """
    file_extension = os.path.splitext(file_path)[1].lower()

    try:
        if file_extension in ['.md', '.txt', '.csv']:
            with open(file_path, 'rb') as file:
                raw = file.read(1024)
                encoding = chardet.detect(raw)['encoding']
            with open(file_path, 'r', encoding=encoding) as file:
                char_count = sum(len(line) for line in file)

        elif file_extension == '.pdf':
            doc = fitz.open(file_path)
            char_count = sum(len(page.get_text()) for page in doc)
            doc.close()

        elif file_extension in ['.jpg', '.png', '.jpeg']:
            # 图片文件无法准确估算字符数，返回True让后续OCR处理
            return None

        elif file_extension == '.docx':
            text = docx2txt.process(file_path)
            char_count = len(text)

        elif file_extension == '.xlsx':
            wb = openpyxl.load_workbook(file_path, read_only=True)
            char_count = sum(len(str(cell.value or '')) for sheet in wb for row in sheet.iter_rows() for cell in row)
            wb.close()

        elif file_extension == '.pptx':
            prs = Presentation(file_path)
            char_count = sum(
                len(shape.text) for slide in prs.slides for shape in slide.shapes if hasattr(shape, 'text'))

        elif file_extension == '.eml':
            with open(file_path, 'r', encoding='utf-8') as file:
                msg = email.message_from_file(file)
                char_count = len(str(msg))

        else:
            # 不支持的文件类型
            return None

        return char_count

    except Exception as e:
        print(f"Error processing file {file_path}: {str(e)}")
        return None


def replace_image_references(text, file_id):
    lines = text.split('\n')
    result = []

    # 匹配带标题的图片引用
    pattern_with_caption = r'^!\[figure\]\((.+\.jpg)\s+(.+)\)$'
    # 匹配不带标题的图片引用
    pattern_without_caption = r'^!\[figure\]\((.+\.jpg)\)$'

    for line in lines:
        if not line.startswith('![figure]'):
            result.append(line)
            continue

        match_with_caption = re.match(pattern_with_caption, line)
        match_without_caption = re.match(pattern_without_caption, line)
        if match_with_caption:
            image_path, caption = match_with_caption.groups()
            debug_logger.info(f"line: {line}, caption: {caption}")
            result.append(f"#### {caption}")
            # result.append(f"![figure](/qanything/assets/file_images/{file_id}/{image_path})")
            result.append(f"![figure]({IMAGES_PROXY_URL}/images/{file_id}/{image_path})")
        elif match_without_caption:
            image_path = match_without_caption.group(1)
            # result.append(f"![figure](/qanything/assets/file_images/{file_id}/{image_path})")
            result.append(f"![figure]({IMAGES_PROXY_URL}/images/{file_id}/{image_path})")
        else:
            result.append(line)

    return '\n'.join(result)




def remove_think_tags(output):
    """
    移除模型输出中的 <think> 标签及其内容。
    
    参数:
        output (str): 模型的原始输出字符串。
        
    返回:
        str: 移除了 <think> 标签及其内容后的字符串。
    """
    # 使用正则表达式匹配 <think> 标签及其内容
    cleaned_output = re.sub(r'<think>.*?</think>', '', output, flags=re.DOTALL)
    
    # 去除多余的空白字符（可选）
    cleaned_output = re.sub(r'\s+', ' ', cleaned_output).strip()
    
    return cleaned_output


def get_docx_toc(docx_path):
    doc = docx.Document(docx_path)
    heading_items = {}
    
    last_level1_heading = None
    last_level2_heading = None
    for para in doc.paragraphs:
        # 检查是否为目录样式或标题样式
        try:
            if para.style.name.startswith('Heading') or para.style.name.startswith('toc'):
                # 只保留到3级标题
                text = para.text.strip().split("\t")[0]  # 只取第一行文本
                if not text:
                    continue  # 跳过空段落
                # print("para.style.name:", para.style.name)
                # print("text:", text)
                if int(para.style.name.split(' ')[1]) <= 3:
                    level = int(para.style.name.split(' ')[1])
                    if level==1 :
                        heading_items[text] = {}
                        last_level1_heading = text
                    elif level==2:
                        if last_level1_heading is None:
                            continue  # 如果没有上一级标题，跳过
                        heading_items[last_level1_heading][text] = []
                        last_level2_heading = text
                    elif level==3:
                        if last_level1_heading is None or last_level2_heading is None:
                            continue  # 如果没有上一级标题，跳过
                        heading_items[last_level1_heading][last_level2_heading].append(text)
                else:
                    continue  # 跳过超过3级以下的标题
            else:
                continue  # 跳过其他样式
        except Exception as e:
            debug_logger.error(f"Error processing paragraph: {str(e)}")
            continue  # 如果处理段落时出错，跳过该段落

        # 暂停一下，需要用户输入enter再继续
        # input("Press Enter to continue...")        
    
    # print("heading_items:", json.dumps(heading_items, indent=4, ensure_ascii=False))
    return str(heading_items)